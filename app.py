from flask import Flask, request, jsonify
import os
import requests
from pptx import Presentation
import io  # Add this line

app = Flask(__name__)

@app.route('/convert', methods=['POST'])
def convert():
    if 'url' not in request.form:
        return jsonify(error="No URL provided"), 400

    url = request.form['url']

    text, error = extract_text_from_url(url)

    if error:
        return jsonify(error=error), 500
    else:
        return jsonify(text=text), 200

def extract_text_from_url(url):
    try:
        response = requests.get(url)
        response.raise_for_status()

        _, extension = os.path.splitext(url)
        extension = extension[1:].lower()  # Removes the dot and makes it lowercase

        if extension == 'pptx':
            prs = Presentation(io.BytesIO(response.content))
            text = "\n".join(
                [slide.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text") for paragraph in shape.text_frame.paragraphs for run in paragraph.runs])
        else:
            return None, "Unsupported filetype"

        return text, None
    except Exception as e:
        return None, str(e)
